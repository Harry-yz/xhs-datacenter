from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x0F)
SURFACE_DARK = RGBColor(0x1A, 0x1A, 0x1A)
SURFACE_ALT = RGBColor(0x24, 0x24, 0x24)
SURFACE_DEEP = RGBColor(0x12, 0x12, 0x12)
ACCENT = RGBColor(0xFF, 0x66, 0x00)
ACCENT_ALT = RGBColor(0xF0, 0x50, 0x00)
ACCENT_LIGHT = RGBColor(0xFF, 0x8C, 0x42)
TEXT_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BODY = RGBColor(0xD1, 0xD1, 0xD1)
TEXT_MUTED = RGBColor(0x99, 0x99, 0x99)
LINE_SOFT = RGBColor(0x3F, 0x3F, 0x3F)

# Fonts (fallback happens in PowerPoint if primary is unavailable)
CN_TITLE_FONT = "MiSans"
EN_TITLE_FONT = "Liter"
BODY_FONT = "Noto Sans SC"


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    left_header: str
    left_points: list[str]
    right_title: str
    right_items: list[str]
    visual_mode: str  # screenshot | abstract
    visual_hint: str
    footnote: str


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_top_brand(slide, page_idx: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.46),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE_DARK
    bar.line.fill.background()

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.34),
        Inches(0.09),
        Inches(3.1),
        Inches(0.28),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = SURFACE_ALT
    tag.line.color.rgb = LINE_SOFT
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "ORAN DATA CENTER"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_LIGHT
    p.font.name = EN_TITLE_FONT
    p.alignment = PP_ALIGN.CENTER

    page = slide.shapes.add_textbox(Inches(12.26), Inches(0.1), Inches(0.75), Inches(0.24))
    p2 = page.text_frame.paragraphs[0]
    p2.text = f"{page_idx:02d}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = EN_TITLE_FONT
    p2.alignment = PP_ALIGN.RIGHT


def add_footer(slide, note: str) -> None:
    foot = slide.shapes.add_textbox(Inches(0.62), Inches(7.02), Inches(12.1), Inches(0.3))
    p = foot.text_frame.paragraphs[0]
    p.text = note
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = BODY_FONT


def add_title_block(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.74), Inches(8.4), Inches(0.95))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(33)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = CN_TITLE_FONT

    subtitle_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.5), Inches(9.5), Inches(0.52))
    ps = subtitle_box.text_frame.paragraphs[0]
    ps.text = subtitle
    ps.font.size = Pt(14)
    ps.font.color.rgb = TEXT_BODY
    ps.font.name = BODY_FONT

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.62),
        Inches(2.03),
        Inches(2.65),
        Inches(0.04),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = ACCENT
    divider.line.fill.background()


def add_left_panel(slide, header: str, points: list[str]) -> None:
    panel_x = 0.62
    panel_y = 2.26
    panel_w = 7.02
    panel_h = 4.58

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(panel_x),
        Inches(panel_y),
        Inches(panel_w),
        Inches(panel_h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE_DARK
    panel.line.color.rgb = LINE_SOFT

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(panel_x + 0.34),
        Inches(panel_y + 0.22),
        Inches(1.72),
        Inches(0.28),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = SURFACE_ALT
    tag.line.fill.background()
    tg = tag.text_frame
    tg.clear()
    p0 = tg.paragraphs[0]
    p0.text = header
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_LIGHT
    p0.font.name = BODY_FONT
    p0.alignment = PP_ALIGN.CENTER

    textbox = slide.shapes.add_textbox(
        Inches(panel_x + 0.34), Inches(panel_y + 0.64), Inches(panel_w - 0.62), Inches(3.8)
    )
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(17 if idx == 0 else 15)
        p.font.color.rgb = TEXT_BODY
        p.font.name = BODY_FONT
        p.space_after = Pt(9)


def add_screenshot_placeholder(slide, x: float, y: float, w: float, h: float, hint: str) -> None:
    outer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    outer.fill.solid()
    outer.fill.fore_color.rgb = SURFACE_ALT
    outer.line.color.rgb = ACCENT_ALT

    inner_margin = 0.18
    inner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + inner_margin),
        Inches(y + inner_margin),
        Inches(w - inner_margin * 2),
        Inches(h - inner_margin * 2),
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = SURFACE_DEEP
    inner.line.color.rgb = LINE_SOFT

    highlight = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x + inner_margin),
        Inches(y + inner_margin),
        Inches(w - inner_margin * 2),
        Inches(0.04),
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ACCENT_LIGHT
    highlight.line.fill.background()

    fade_top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x + inner_margin),
        Inches(y + inner_margin + 0.04),
        Inches(w - inner_margin * 2),
        Inches(0.46),
    )
    fade_top.fill.solid()
    fade_top.fill.fore_color.rgb = SURFACE_DARK
    fade_top.fill.transparency = 0.35
    fade_top.line.fill.background()

    fade_bottom = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x + inner_margin),
        Inches(y + h - inner_margin - 0.54),
        Inches(w - inner_margin * 2),
        Inches(0.54),
    )
    fade_bottom.fill.solid()
    fade_bottom.fill.fore_color.rgb = SURFACE_DARK
    fade_bottom.fill.transparency = 0.2
    fade_bottom.line.fill.background()

    center_box = slide.shapes.add_textbox(
        Inches(x + 0.36), Inches(y + h / 2 - 0.35), Inches(w - 0.72), Inches(0.95)
    )
    tf = center_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"【请替换：{hint}】"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_BODY
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "建议直接放白底产品截图，容器已做融合处理"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = BODY_FONT
    p2.alignment = PP_ALIGN.CENTER


def add_abstract_visual(slide, x: float, y: float, w: float, h: float, label: str) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE_ALT
    panel.line.color.rgb = LINE_SOFT

    frame_count = 4
    for i in range(frame_count):
        pad = 0.28 + i * 0.2
        frame = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + pad),
            Inches(y + 0.42 + pad * 0.38),
            Inches(w - 2 * pad),
            Inches(h - 1.05 - pad * 0.68),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = SURFACE_DARK
        frame.fill.transparency = 0.25 + i * 0.14
        frame.line.color.rgb = ACCENT if i % 2 == 0 else LINE_SOFT

    axis = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x + 0.38),
        Inches(y + h - 0.68),
        Inches(w - 0.76),
        Inches(0.03),
    )
    axis.fill.solid()
    axis.fill.fore_color.rgb = ACCENT
    axis.line.fill.background()

    line1 = slide.shapes.add_connector(
        1, Inches(x + 0.72), Inches(y + h - 1.0), Inches(x + w - 0.72), Inches(y + 0.95)
    )
    line1.line.color.rgb = ACCENT_LIGHT
    line1.line.width = Pt(1.2)

    line2 = slide.shapes.add_connector(
        1, Inches(x + 1.2), Inches(y + h - 1.15), Inches(x + w - 1.15), Inches(y + 1.32)
    )
    line2.line.color.rgb = LINE_SOFT
    line2.line.width = Pt(0.8)

    label_box = slide.shapes.add_textbox(Inches(x + 0.34), Inches(y + 0.2), Inches(w - 0.68), Inches(0.38))
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_LIGHT
    p.font.name = EN_TITLE_FONT

    chips = ["STRUCTURED", "PRECISE", "FORWARD"]
    cx = x + 0.44
    for chip in chips:
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(cx),
            Inches(y + h - 0.52),
            Inches(1.35),
            Inches(0.28),
        )
        c.fill.solid()
        c.fill.fore_color.rgb = SURFACE_DARK
        c.line.color.rgb = LINE_SOFT
        cp = c.text_frame.paragraphs[0]
        cp.text = chip
        cp.font.size = Pt(8.5)
        cp.font.color.rgb = TEXT_MUTED
        cp.font.name = EN_TITLE_FONT
        cp.alignment = PP_ALIGN.CENTER
        cx += 1.43


def add_right_panel(slide, spec: SlideSpec) -> None:
    panel_x = 7.74
    panel_y = 2.26
    panel_w = 4.98
    panel_h = 4.58

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(panel_x),
        Inches(panel_y),
        Inches(panel_w),
        Inches(panel_h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE_ALT
    panel.line.color.rgb = LINE_SOFT

    title_box = slide.shapes.add_textbox(Inches(panel_x + 0.28), Inches(panel_y + 0.2), Inches(panel_w - 0.56), Inches(0.32))
    t = title_box.text_frame.paragraphs[0]
    t.text = spec.right_title
    t.font.size = Pt(14)
    t.font.bold = True
    t.font.color.rgb = ACCENT_LIGHT
    t.font.name = CN_TITLE_FONT

    y = panel_y + 0.56
    for item in spec.right_items:
        row = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(panel_x + 0.24),
            Inches(y),
            Inches(panel_w - 0.48),
            Inches(0.43),
        )
        row.fill.solid()
        row.fill.fore_color.rgb = SURFACE_DARK
        row.line.color.rgb = LINE_SOFT
        rp = row.text_frame.paragraphs[0]
        rp.text = item
        rp.font.size = Pt(12)
        rp.font.color.rgb = TEXT_BODY
        rp.font.name = BODY_FONT
        y += 0.51

    vis_x = panel_x + 0.24
    vis_w = panel_w - 0.48
    vis_top = y + 0.08
    vis_h = panel_y + panel_h - vis_top - 0.2
    if vis_h < 1.54:
        vis_top = panel_y + 2.74
        vis_h = 1.62

    if spec.visual_mode == "screenshot":
        add_screenshot_placeholder(slide, vis_x, vis_top, vis_w, vis_h, spec.visual_hint)
    else:
        add_abstract_visual(slide, vis_x, vis_top, vis_w, vis_h, spec.visual_hint)


def add_content_slide(prs: Presentation, page_idx: int, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_brand(slide, page_idx)
    add_title_block(slide, spec.title, spec.subtitle)
    add_left_panel(slide, spec.left_header, spec.left_points)
    add_right_panel(slide, spec)
    add_footer(slide, spec.footnote)


def add_cover_slide(prs: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_brand(slide, 1)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.36), Inches(11.1), Inches(1.35))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = CN_TITLE_FONT

    sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(2.84), Inches(10.6), Inches(0.68))
    ps = sub_box.text_frame.paragraphs[0]
    ps.text = subtitle
    ps.font.size = Pt(20)
    ps.font.color.rgb = TEXT_BODY
    ps.font.name = BODY_FONT

    accent_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.68),
        Inches(3.74),
        Inches(8.15),
        Inches(0.09),
    )
    accent_band.fill.solid()
    accent_band.fill.fore_color.rgb = ACCENT
    accent_band.line.fill.background()

    chips = [
        "已跑通最小闭环",
        "XHS 主线可用",
        "TikTok 并行接入",
        datetime.now().strftime("阶段同步版 · %Y-%m-%d"),
    ]
    x = 0.68
    y = 4.18
    widths = [2.2, 1.9, 2.1, 2.55]
    for i, chip in enumerate(chips):
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(widths[i]),
            Inches(0.44),
        )
        c.fill.solid()
        c.fill.fore_color.rgb = SURFACE_ALT
        c.line.color.rgb = LINE_SOFT
        cp = c.text_frame.paragraphs[0]
        cp.text = chip
        cp.font.size = Pt(11)
        cp.font.color.rgb = ACCENT_LIGHT if i < 3 else TEXT_MUTED
        cp.font.bold = i < 3
        cp.font.name = BODY_FONT
        cp.alignment = PP_ALIGN.CENTER
        x += widths[i] + 0.2

    add_abstract_visual(slide, 8.54, 2.34, 4.18, 3.65, "CLOSED LOOP")
    add_footer(slide, footer)


def build_internal_v3(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    common_note = "阶段同步口径：当前以线上可见能力和链路稳定性为主，量化KPI待测试体系完善后补充。"

    add_cover_slide(
        prs,
        title="Oran 数据中台阶段汇报",
        subtitle="已跑通最小闭环，进入稳定迭代阶段（XHS主线 + TikTok并行）",
        footer=common_note,
    )

    slides = [
        SlideSpec(
            title="执行摘要",
            subtitle="这阶段核心不是做大而是做通：从概念验证进入可用阶段",
            left_header="本阶段结论",
            left_points=[
                "结论：项目已从“能演示”进入“可使用”，最小闭环具备持续迭代基础。",
                "完成项：入口-二级页-搜索页链路已打通，核心交互和登录拦截稳定。",
                "成熟度：当前处于V1可用阶段，优先做稳命中、补采和解释能力。",
                "本次目标：阶段同步，不做超前承诺，明确下一阶段边界与节奏。",
            ],
            right_title="阶段价值",
            right_items=[
                "业务侧可直接演示完整路径",
                "产品侧具备可迭代骨架",
                "技术侧明确优化主线",
                "跨平台扩展具备同构基础",
            ],
            visual_mode="abstract",
            visual_hint="EXEC SUMMARY",
            footnote=common_note,
        ),
        SlideSpec(
            title="已上线能力全景",
            subtitle="从入口到洞察，核心链路已连续可用",
            left_header="闭环能力",
            left_points=[
                "首页：六平台入口和统一视觉入口已完成，可承载平台化叙事。",
                "XHS二级页：总量、趋势、行业入口协同，支持从看板到定位问题。",
                "搜索三级页：品类/达人检索、排序分页、外链查看链路可跑通。",
                "体验层：中英文、主题切换、登录拦截与会话衔接均已落位。",
            ],
            right_title="链路视图",
            right_items=[
                "入口层：平台分发",
                "分析层：概览与趋势",
                "检索层：结果与筛选",
                "动作层：详情外链与回看",
            ],
            visual_mode="screenshot",
            visual_hint="首页 + 二级页 + 三级页拼图",
            footnote=common_note,
        ),
        SlideSpec(
            title="可演示能力（能看 / 能搜 / 能定位）",
            subtitle="围绕真实使用动作组织能力，而非功能堆砌",
            left_header="演示要点",
            left_points=[
                "能看：二级页提供总量和趋势窗口，可快速建立盘面认知。",
                "能搜：三级页支持关键词与达人双入口，排序与分页流程顺畅。",
                "能定位：行业入口可直达目标搜索，减少反复跳转成本。",
                "能回链：结果支持外链原文，便于业务和运营做二次判断。",
            ],
            right_title="讲解提示",
            right_items=[
                "先演示二级页总览",
                "再演示搜索筛选过程",
                "最后展示结果与外链",
                "建议控制在45秒内",
            ],
            visual_mode="screenshot",
            visual_hint="二级页 + 搜索页关键画面",
            footnote=common_note,
        ),
        SlideSpec(
            title="当前不足与根因",
            subtitle="问题已识别且可拆解，下一步按优先级逐项收敛",
            left_header="主要问题",
            left_points=[
                "命中稳定性不足：大词、长尾词在不同时间窗口表现波动。",
                "字段完整性不均：互动和达人相关字段存在批次差异。",
                "性能体验波动：冷启动与重查询场景下等待时间偏长。",
                "解释能力薄弱：结果为什么出现、数据新鲜度展示不够直观。",
            ],
            right_title="根因拆解",
            right_items=[
                "召回策略覆盖不足",
                "补采队列受限额影响",
                "字段来源链路不均衡",
                "缓存与SQL策略待优化",
            ],
            visual_mode="abstract",
            visual_hint="ROOT CAUSE MAP",
            footnote=common_note,
        ),
        SlideSpec(
            title="下一阶段主线A（XHS深耕）",
            subtitle="先把命中和补采做稳，再补齐可解释层",
            left_header="优先动作",
            left_points=[
                "命中增强：扩大关键词裂变与宽匹配覆盖，减少空结果场景。",
                "补采闭环：低命中或老化触发补采，回填后再读取主结果。",
                "并发治理：同词请求合并，避免重复调用上游接口。",
                "解释层补齐：在结果页显示命中说明与数据新鲜度提示。",
            ],
            right_title="阶段产出",
            right_items=[
                "空结果比例可见下降",
                "字段缺失率逐步收敛",
                "解释信息可被业务读懂",
                "接口成本可持续控制",
            ],
            visual_mode="screenshot",
            visual_hint="搜索结果页命中与解释层截图",
            footnote=common_note,
        ),
        SlideSpec(
            title="TikTok并行策略",
            subtitle="按同构原则并行接入，不拖慢XHS主线节奏",
            left_header="并行原则",
            left_points=[
                "字段同构优先：内容、达人、互动、发布时间等先对齐核心字段。",
                "页面最小闭环：先提供可进入的轻量总览与基础搜索能力。",
                "交互模型复用：尽量沿用XHS交互，降低学习和维护成本。",
                "迭代策略：先可用再做深，不一次性扩张全部分析模块。",
            ],
            right_title="分阶段推进",
            right_items=[
                "Phase 1：接入与映射",
                "Phase 2：总览与检索",
                "Phase 3：品牌达人深化",
                "Phase 4：与XHS统一运营视角",
            ],
            visual_mode="screenshot",
            visual_hint="TikTok轻量总览与搜索页截图",
            footnote=common_note,
        ),
        SlideSpec(
            title="里程碑与协同需求",
            subtitle="按M1/M2/M3推进，阶段目标可检视、可复盘",
            left_header="阶段节奏",
            left_points=[
                "M1：命中优化与数据质量面板上线，先确保基础稳定可解释。",
                "M2：品牌/达人分析增强上线，提升业务侧分析效率。",
                "M3：TikTok轻量能力上线，形成跨平台并行能力雏形。",
                "协同点：接口配额、字段口径、联调窗口要提前锁定。",
            ],
            right_title="协同关注",
            right_items=[
                "接口配额与上游节奏",
                "字段映射与口径一致",
                "联调窗口与责任人",
                "版本验收与回归机制",
            ],
            visual_mode="screenshot",
            visual_hint="里程碑看板或联调排期截图",
            footnote=common_note,
        ),
        SlideSpec(
            title="收尾结论",
            subtitle="本阶段先证明“做得通”，下一阶段聚焦“做得稳、做得深”",
            left_header="一句话总结",
            left_points=[
                "我们已经完成最小闭环验证：链路可演示、能力可复用、方向可持续。",
                "下一阶段坚持聚焦：XHS主线做深，TikTok并行做轻，不盲目扩张。",
                "节奏原则不变：先稳定再增强，先可解释再规模化。",
                "本次同步的价值是统一认知和节奏，为后续协同减少偏差。",
            ],
            right_title="下阶段协同节奏",
            right_items=[
                "每周同步：问题与进展",
                "双周复盘：阶段目标达成",
                "关键变更：先对齐再执行",
                "对外口径：以已上线能力为准",
            ],
            visual_mode="abstract",
            visual_hint="NEXT STEP",
            footnote=common_note,
        ),
    ]

    for idx, spec in enumerate(slides, start=2):
        add_content_slide(prs, idx, spec)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def write_speaker_notes_v3(path: Path) -> None:
    notes = """# Oran 数据中台内部汇报讲稿 v3（5-8分钟）

## 第1页 封面（30秒）
- 我们这次汇报聚焦一个结论：数据中台已经跑通最小闭环，项目从概念验证进入可用阶段。
- 当前策略是 XHS 主线持续做深，TikTok 以同构方式并行接入，目标是稳步扩展而不是一次铺开。
- 这次是阶段同步会，不追求一次性讲完所有规划，重点讲清楚“做到了什么、下一步做什么”。

## 第2页 执行摘要（45秒）
- 先说结论：这阶段最重要的成果不是功能数量，而是核心链路可用，已经具备持续迭代基础。
- 从入口到二级页再到搜索页，关键页面和核心交互都已经能完整演示，团队协同也有了统一骨架。
- 当前成熟度判断是 V1 可用，我们下一步会把命中稳定性、补采闭环和解释层作为优先主线推进。
- 所以这次汇报的目标是对齐阶段认知，不做超前承诺，把节奏和边界讲清楚。

## 第3页 已上线能力全景（45秒）
- 现在可以把它当成一条完整链路来看：首页负责平台入口，二级页负责总览和趋势，三级页负责检索和定位。
- 我们不是只做了一个展示页，而是打通了“看-搜-定位-回链”的操作路径，业务方可以按真实动作使用。
- 体验层上，中英文、主题切换、登录拦截和会话衔接都已经落位，保证使用过程不割裂。
- 右侧这块你可以放链路拼图截图，讲的时候按入口到搜索的顺序扫一遍就很顺。

## 第4页 可演示能力（45秒）
- 这一页建议按“能看、能搜、能定位”来讲，听众会比按功能列表更容易理解价值。
- 能看，是二级页建立盘面；能搜，是三级页完成检索和筛选；能定位，是行业入口和外链闭环能支持具体判断。
- 重点不是讲每个按钮，而是证明这套路径已经可以支撑一次完整分析动作。
- 右侧建议放二级页和搜索页关键截图，配合你的现场操作说明。

## 第5页 当前不足与根因（50秒）
- 当前问题我们不回避，主要集中在命中稳定性、字段完整性、性能波动和结果解释性这四类。
- 这些问题并不是方向错误，而是典型的 V1 到 V1.5 过渡问题，本质上是策略、队列和查询治理还不够成熟。
- 我们已经把根因拆成可执行模块，不会做泛泛而谈的优化，而是按优先级逐项收敛。
- 这页讲完，听众通常会更关心下一步怎么落地，所以自然过渡到主线A。

## 第6页 下一阶段主线A（50秒）
- 下一阶段我们聚焦四件事：命中增强、补采闭环、并发治理、解释层补齐。
- 命中增强解决“搜不出来”，补采闭环解决“数据不新”，并发治理解决“成本和稳定性”，解释层解决“业务看不懂”。
- 这四件事放在一起，目标是让结果不仅出现，而且可信、可读、可复用。
- 右侧建议放搜索结果页与解释层截图，这会让改进方向更直观。

## 第7页 TikTok并行策略（45秒）
- TikTok 并行不是另起炉灶，而是按同构逻辑复用 XHS 已有模型，避免把团队节奏拉散。
- 我们先做最小闭环，再做深度分析，这样能在控制风险的同时保留扩展速度。
- 节奏上按分阶段推进：先接入映射，再总览检索，最后再做品牌达人深化。
- 右侧建议放 TikTok 轻量总览或搜索页截图，强化“并行但可控”的感受。

## 第8页 里程碑与协同需求（45秒）
- 里程碑按 M1、M2、M3 推进，核心是每一段都有可检视的阶段结果，不做泛目标。
- 对团队协同的要求主要是三块：接口配额、字段口径、联调窗口，这三块提前锁住可以显著降低返工。
- 这页不需要讲太细的项目管理术语，讲清楚“为什么这些协同是关键阻断点”就够了。
- 右侧可以放排期或里程碑看板截图，帮助大家快速形成共同时间认知。

## 第9页 收尾结论（35秒）
- 最后总结一句：我们已经证明这件事做得通，下一阶段要把它做得稳、做得深。
- 节奏上继续坚持先稳定再增强、先可解释再规模化，确保每一步都能沉淀成可复用能力。
- 本次阶段同步的价值是统一认知和节奏，为后续协同减少偏差，而不是在会上一次性拍板所有资源。

---

## 截图替换说明（你后续手动加入）
- 第3页：`首页 + XHS二级页 + 搜索三级页` 链路拼图
- 第4页：`二级页总览 + 搜索筛选结果` 关键演示画面
- 第6页：`搜索结果页命中提升 + 解释层` 画面
- 第7页：`TikTok轻量总览或搜索页` 画面
- 第8页：`里程碑看板 / 联调排期` 画面
- 建议优先使用 16:9 截图，直接覆盖占位容器即可；容器已做白底融合效果，无需二次修图。
"""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(notes, encoding="utf-8")


def main() -> None:
    base = Path("/opt/xhs_data_center/docs/ppt")
    internal_ppt = base / "oran-datacenter-internal-v3.pptx"
    notes_file = base / "oran-datacenter-speaker-notes-v3.md"

    build_internal_v3(internal_ppt)
    write_speaker_notes_v3(notes_file)

    print(f"generated: {internal_ppt}")
    print(f"generated: {notes_file}")


if __name__ == "__main__":
    main()
